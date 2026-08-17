"""Stage 1: incremental detection and text extraction."""

from __future__ import annotations

from pathlib import Path

import pytest

from pipeline.config import Settings
from pipeline.ingest import MAX_CHARS, ingest_sources, mark_processed
from pipeline.state import RunStore
from pipeline.storage import Storage


@pytest.fixture
def workspace(tmp_path: Path) -> tuple[Settings, RunStore]:
    """Settings rooted at a temp directory, so tests never touch the real repo.

    Passing `repo_root` explicitly rather than monkeypatching: an earlier version patched
    the property and had no effect, because get_settings() is cached — so the tests wrote
    into the production sources/ directory and overwrote its README. `repo_root` is a real
    field now precisely so this is impossible.
    """
    (tmp_path / "sources").mkdir()
    (tmp_path / "state").mkdir()

    settings = Settings(repo_root=tmp_path)
    return settings, RunStore(settings)


def write(settings: Settings, name: str, content: str = "Some project content.") -> Path:
    path = settings.sources_dir / name
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")
    return path


# --- incremental detection ---------------------------------------------------


def test_new_file_is_found(workspace: tuple[Settings, RunStore]) -> None:
    settings, store = workspace
    write(settings, "prd.md")
    sources = ingest_sources(store, settings)
    assert [s.name for s in sources] == ["prd.md"]
    assert sources[0].text == "Some project content."


def test_processed_file_is_skipped(workspace: tuple[Settings, RunStore]) -> None:
    settings, store = workspace
    write(settings, "prd.md")
    sources = ingest_sources(store, settings)
    mark_processed(sources, store, settings)
    assert ingest_sources(store, settings) == []


def test_changed_file_is_reprocessed(workspace: tuple[Settings, RunStore]) -> None:
    settings, store = workspace
    write(settings, "prd.md", "v1")
    mark_processed(ingest_sources(store, settings), store, settings)

    write(settings, "prd.md", "v2 with new results")
    sources = ingest_sources(store, settings)
    assert len(sources) == 1
    assert sources[0].text == "v2 with new results"


def test_detection_is_by_content_not_mtime(workspace: tuple[Settings, RunStore]) -> None:
    """A git clone rewrites every mtime.

    Keying on mtime would make the first run in a fresh container re-extract the whole
    archive and file duplicate review rows for work already approved.
    """
    settings, store = workspace
    path = write(settings, "prd.md", "stable content")
    mark_processed(ingest_sources(store, settings), store, settings)

    # Same bytes, much later mtime — as a fresh checkout would produce.
    import os

    os.utime(path, (10**9, 10**9))
    assert ingest_sources(store, settings) == []


def test_unprocessed_on_failure(workspace: tuple[Settings, RunStore]) -> None:
    """The index only advances via mark_processed.

    So a crash between extraction and reconcile leaves the file pending rather than
    silently skipped on the next run.
    """
    settings, store = workspace
    write(settings, "prd.md")
    ingest_sources(store, settings)  # discovered, but never marked
    assert len(ingest_sources(store, settings)) == 1


# --- filtering --------------------------------------------------------------


def test_readme_is_not_a_source(workspace: tuple[Settings, RunStore]) -> None:
    """sources/README.md documents the folder; it is not material to extract from."""
    settings, store = workspace
    write(settings, "README.md", "# how to use this folder")
    write(settings, "prd.md")
    assert [s.name for s in ingest_sources(store, settings)] == ["prd.md"]


def test_unsupported_extensions_ignored(workspace: tuple[Settings, RunStore]) -> None:
    settings, store = workspace
    write(settings, "notes.md")
    write(settings, "screenshot.png", "not really an image")
    write(settings, "data.xlsx", "not really a spreadsheet")
    assert [s.name for s in ingest_sources(store, settings)] == ["notes.md"]


def test_hidden_files_ignored(workspace: tuple[Settings, RunStore]) -> None:
    settings, store = workspace
    write(settings, ".DS_Store", "junk")
    write(settings, ".gitkeep", "")
    assert ingest_sources(store, settings) == []


def test_empty_file_skipped(workspace: tuple[Settings, RunStore]) -> None:
    """Nothing to extract, and an empty prompt would waste a call."""
    settings, store = workspace
    write(settings, "blank.md", "   \n\n  ")
    assert ingest_sources(store, settings) == []


def test_nested_directories_are_searched(workspace: tuple[Settings, RunStore]) -> None:
    """You'll organize sources into folders; the walk has to follow."""
    settings, store = workspace
    write(settings, "2026/cathay/prd.md")
    assert [s.name for s in ingest_sources(store, settings)] == ["prd.md"]


def test_oversized_text_is_truncated(workspace: tuple[Settings, RunStore]) -> None:
    """One mis-saved file must not consume a run's whole token budget."""
    settings, store = workspace
    write(settings, "huge.md", "x" * (MAX_CHARS + 5000))
    sources = ingest_sources(store, settings)
    assert len(sources[0].text or "") == MAX_CHARS


# --- format handling --------------------------------------------------------


def test_pdf_is_passed_as_bytes(workspace: tuple[Settings, RunStore]) -> None:
    """PDFs go to the model natively: layout is where a report's structure lives, and
    Claude reads it better than a local text extractor."""
    settings, store = workspace
    (settings.sources_dir / "paper.pdf").write_bytes(b"%PDF-1.4 fake content")
    sources = ingest_sources(store, settings)
    assert sources[0].is_pdf
    assert sources[0].text is None
    assert sources[0].pdf_bytes is not None


def test_oversized_pdf_is_skipped(workspace: tuple[Settings, RunStore]) -> None:
    """Anthropic caps a request at 32MB; a too-large PDF would fail the call."""
    settings, store = workspace
    from pipeline.ingest import MAX_PDF_BYTES

    (settings.sources_dir / "huge.pdf").write_bytes(b"%PDF" + b"x" * MAX_PDF_BYTES)
    assert ingest_sources(store, settings) == []


def test_docx_text_includes_table_cells(workspace: tuple[Settings, RunStore]) -> None:
    """A report's deliverables are as likely to be in a table as in prose."""
    from docx import Document

    settings, store = workspace
    document = Document()
    document.add_paragraph("Project summary paragraph.")
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Deliverable"
    table.rows[0].cells[1].text = "Shipped the BRD agent"
    document.save(str(settings.sources_dir / "report.docx"))

    text = ingest_sources(store, settings)[0].text or ""
    assert "Project summary paragraph." in text
    assert "Shipped the BRD agent" in text


def test_pptx_text_labels_slides(workspace: tuple[Settings, RunStore]) -> None:
    """Slide numbers give the model structure: an agenda slide reads differently from
    a results slide."""
    from pptx import Presentation
    from pptx.util import Inches

    settings, store = workspace
    presentation = Presentation()
    for content in ("Agenda", "Results: 40% faster"):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.text = content
    presentation.save(str(settings.sources_dir / "deck.pptx"))

    text = ingest_sources(store, settings)[0].text or ""
    assert "--- Slide 1 ---" in text
    assert "--- Slide 2 ---" in text
    assert "Results: 40% faster" in text


def test_unreadable_file_does_not_fail_the_run(workspace: tuple[Settings, RunStore]) -> None:
    """A corrupt .docx must not stop a resume rebuild from already-approved rows."""
    settings, store = workspace
    (settings.sources_dir / "corrupt.docx").write_bytes(b"this is not a zip archive")
    write(settings, "good.md")
    assert [s.name for s in ingest_sources(store, settings)] == ["good.md"]


def test_missing_sources_directory_is_not_an_error(tmp_path: Path) -> None:
    settings = Settings(repo_root=tmp_path)
    assert ingest_sources(RunStore(settings), settings) == []


# --- the intake folder must arrive from durable storage -----------------------


class _DurableOnlyStorage(Storage):
    """Storage with no filesystem behind it, like the GitHub backend on Cloud Run.

    LocalStorage cannot express this case: there, `sources/` on disk *is* the durable copy,
    so a test using it passes whether or not ingest fetches anything. That is exactly why
    the production gap — an empty `sources/` in the container — went unnoticed.
    """

    def __init__(self, data: dict[str, bytes]) -> None:
        self._data = dict(data)

    def read(self, path: str) -> bytes | None:
        return self._data.get(path)

    def write(self, path: str, data: bytes, message: str) -> None:  # noqa: ARG002
        self._data[path] = data

    def delete_prefix(self, prefix: str, message: str) -> int:  # noqa: ARG002
        keys = [k for k in self._data if k.startswith(prefix)]
        for key in keys:
            del self._data[key]
        return len(keys)

    def list_prefix(self, prefix: str) -> list[str]:
        return sorted({k for k in self._data if k.startswith(prefix)})

    def walk(self, prefix: str) -> list[str]:
        return sorted(k for k in self._data if k.startswith(prefix))


def test_sources_pushed_to_the_repo_are_ingested_on_an_empty_container(tmp_path: Path) -> None:
    """A document pushed to `sources/` must be picked up by a container that has never seen it.

    The measured gap: the Dockerfile creates an empty `/app/sources`, the entrypoint skips the
    git checkout when STORAGE_BACKEND=github, and ingest scanned that empty directory. So
    dropping a file in `sources/` and pushing it did nothing, with no log line saying so.
    """
    (tmp_path / "sources").mkdir()
    (tmp_path / "state").mkdir()
    settings = Settings(repo_root=tmp_path)
    store = RunStore(
        settings,
        storage=_DurableOnlyStorage(
            {
                "sources/internship_report.md": b"Built an ETL pipeline on AWS Lambda.",
                "sources/notes/extra.txt": b"Presented at a conference.",
                "sources/README.md": b"instructions, not material",
                "sources/.gitkeep": b"",
            }
        ),
    )

    assert not any(settings.sources_dir.rglob("*")), "the container starts with nothing"

    pending = ingest_sources(store, settings)
    names = sorted(s.name for s in pending)

    assert names == ["extra.txt", "internship_report.md"], (
        f"expected both pushed documents to be ingested, got {names}"
    )
    # README.md and dotfiles are the folder's own scaffolding, not source material.
    assert "README.md" not in names


def test_materialize_sources_does_not_rewrite_unchanged_files(tmp_path: Path) -> None:
    """Called on every run, so it must be a no-op when nothing changed — otherwise the
    local backend would rewrite its own files each time."""
    (tmp_path / "sources").mkdir()
    (tmp_path / "state").mkdir()
    settings = Settings(repo_root=tmp_path)
    store = RunStore(settings, storage=_DurableOnlyStorage({"sources/a.md": b"same bytes"}))

    assert store.materialize_sources() == 1
    assert store.materialize_sources() == 0


def test_changed_source_is_refetched_on_a_reused_instance(tmp_path: Path) -> None:
    """A stale local copy must not shadow an updated document."""
    (tmp_path / "sources").mkdir()
    (tmp_path / "state").mkdir()
    settings = Settings(repo_root=tmp_path)
    storage = _DurableOnlyStorage({"sources/a.md": b"version one"})
    store = RunStore(settings, storage=storage)

    store.materialize_sources()
    storage.write("sources/a.md", b"version two, revised", "update")

    assert store.materialize_sources() == 1
    assert (settings.sources_dir / "a.md").read_bytes() == b"version two, revised"


# --- HTML sources -------------------------------------------------------------
#
# Saved web pages are the messiest thing this folder accepts: "Save as HTML" from Word or a
# browser carries minified JavaScript, CSS, vendor namespaces (<o:p>) and conditional
# comments. Everything below is about not feeding that to the model, which would happily
# invent a "project" out of variable names.

SAVED_PAGE = """<!DOCTYPE html>
<html><head><title>Internship Closing Report</title>
<style>.hdr{color:red}</style>
<script>var tracking={id:"abc"};function boot(){console.log("delivered")}</script>
</head><body>
<h1>AI News Agent &amp; RAG Pipeline</h1>
<p>Built a routing chatbot.</p>
<p>Cut ETL runtime by 40%.</p>
<table>
<tr><th>Deliverable</th><th>Status</th></tr>
<tr><td>BRD agent</td><td>Shipped &lt;v1.0&gt;</td></tr>
</table>
<noscript>Enable JavaScript</noscript>
</body></html>"""


def test_html_is_ingested(workspace: tuple[Settings, RunStore]) -> None:
    settings, store = workspace
    write(settings, "report.html", SAVED_PAGE)

    sources = ingest_sources(store, settings)

    assert [s.name for s in sources] == ["report.html"]
    assert "Built a routing chatbot." in (sources[0].text or "")


def test_htm_extension_also_works(workspace: tuple[Settings, RunStore]) -> None:
    """Windows "Save as" still produces .htm."""
    settings, store = workspace
    write(settings, "old.htm", "<html><body><p>Ran a pilot study.</p></body></html>")

    assert [s.name for s in ingest_sources(store, settings)] == ["old.htm"]


def test_script_and_style_content_never_reaches_the_model(
    workspace: tuple[Settings, RunStore],
) -> None:
    """The single most important property of HTML extraction."""
    settings, store = workspace
    write(settings, "report.html", SAVED_PAGE)

    text = ingest_sources(store, settings)[0].text or ""

    for leaked in ("console.log", "var tracking", "color:red", "boot()"):
        assert leaked not in text, f"{leaked!r} leaked into the extracted text"
    # <noscript> is fallback copy for a browser, not content about the author.
    assert "Enable JavaScript" not in text


def test_block_tags_keep_sentences_apart(workspace: tuple[Settings, RunStore]) -> None:
    """Without newlines on block tags, two unrelated sentences fuse into one nonsense
    phrase — worse than losing either, because the model treats it as a single claim.

    The markup here has NO whitespace between the tags, which is what Word exports and
    minified pages actually look like. An earlier version of this test used a document with
    newlines between the `<p>` elements, so the source whitespace did the separating and the
    test passed even with block-tag handling removed entirely — verified by removing it.
    """
    settings, store = workspace
    write(
        settings,
        "minified.html",
        "<html><body><p>Built a routing chatbot.</p><p>Cut ETL runtime by 40%.</p>"
        "<div>Owned the deploy.</div><ul><li>First</li><li>Second</li></ul></body></html>",
    )

    text = ingest_sources(store, settings)[0].text or ""

    assert "chatbot.Cut" not in text, "adjacent <p> blocks fused into one sentence"
    assert "40%.Owned" not in text
    assert "FirstSecond" not in text, "adjacent <li> items fused"
    assert "Built a routing chatbot." in text
    assert "Cut ETL runtime by 40%." in text


def test_entities_are_decoded(workspace: tuple[Settings, RunStore]) -> None:
    settings, store = workspace
    write(settings, "report.html", SAVED_PAGE)

    text = ingest_sources(store, settings)[0].text or ""

    assert "AI News Agent & RAG Pipeline" in text
    assert "Shipped <v1.0>" in text
    assert "&amp;" not in text


def test_table_cells_are_joined_like_docx(workspace: tuple[Settings, RunStore]) -> None:
    """A deliverables table is often where the concrete detail lives, and a row reads as one
    record. Same separator as `_docx_text` so both sources look alike to the model."""
    settings, store = workspace
    write(settings, "report.html", SAVED_PAGE)

    text = ingest_sources(store, settings)[0].text or ""

    assert "BRD agent | Shipped <v1.0>" in text


def test_html_with_no_prose_is_skipped(workspace: tuple[Settings, RunStore]) -> None:
    """A page that is only markup and scripts yields nothing, and must be skipped rather
    than filed as an empty source."""
    settings, store = workspace
    write(
        settings,
        "empty.html",
        "<html><head><script>var a=1</script><style>p{}</style></head><body></body></html>",
    )

    assert ingest_sources(store, settings) == []


def test_nested_script_does_not_swallow_the_rest_of_the_page(
    workspace: tuple[Settings, RunStore],
) -> None:
    """A depth counter, not a boolean: with a flag, the first `</script>` flips skipping off
    (or a stray close flips it on) and the remainder of the document is lost or leaked."""
    settings, store = workspace
    write(
        settings,
        "weird.html",
        "<html><body><script>var a=1</script>"
        "<p>Led the migration.</p>"
        "<script>var b=2</script>"
        "<p>Wrote the runbook.</p></body></html>",
    )

    text = ingest_sources(store, settings)[0].text or ""

    assert "Led the migration." in text
    assert "Wrote the runbook." in text
    assert "var a=1" not in text
    assert "var b=2" not in text


def test_malformed_html_still_yields_text(workspace: tuple[Settings, RunStore]) -> None:
    """Word and Confluence exports are not well-formed. Unclosed tags and vendor namespaces
    must degrade to "some text extracted", never to an exception."""
    settings, store = workspace
    write(
        settings,
        "word.html",
        "<html><body><p>Shipped the pilot"
        "<div><o:p>Delivered on time</o:p>"
        "<!--[if gte mso 9]><xml><o:Settings/></xml><![endif]-->"
        "<b>Unclosed bold</body>",
    )

    text = ingest_sources(store, settings)[0].text or ""

    assert "Shipped the pilot" in text
    assert "Delivered on time" in text
    assert "o:Settings" not in text
