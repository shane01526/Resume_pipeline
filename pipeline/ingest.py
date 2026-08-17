"""Stage 1: find source files worth extracting from.

Incremental by content hash, not mtime: a `git clone` in a fresh container rewrites every
mtime, which would make the first run re-extract the entire archive and file duplicate
review rows for work you already approved.

Text extraction happens here rather than in stage 2 so the LLM call receives plain text
for everything except PDFs, which go to the model natively — Claude reads a PDF's layout
better than any local text extractor, and layout is where a project report's structure
lives.
"""

from __future__ import annotations

import hashlib
import logging
from dataclasses import dataclass
from html.parser import HTMLParser
from pathlib import Path

from pipeline.config import Settings
from pipeline.state import RunStore

log = logging.getLogger(__name__)

TEXT_SUFFIXES = frozenset({".md", ".txt", ".markdown", ".rst"})
DOCX_SUFFIXES = frozenset({".docx"})
PPTX_SUFFIXES = frozenset({".pptx"})
PDF_SUFFIXES = frozenset({".pdf"})
HTML_SUFFIXES = frozenset({".html", ".htm"})

SUPPORTED = TEXT_SUFFIXES | DOCX_SUFFIXES | PPTX_SUFFIXES | PDF_SUFFIXES | HTML_SUFFIXES

# A single source over ~200k characters is almost certainly a book, a transcript dump, or
# a mis-saved binary. Truncating keeps one bad file from consuming a run's whole budget.
MAX_CHARS = 200_000

# Anthropic's PDF limit is 32MB per request; stay well under it.
MAX_PDF_BYTES = 20 * 1024 * 1024


@dataclass(slots=True)
class Source:
    """One file to extract from."""

    path: Path
    sha256: str
    #: Plain text for md/txt/docx/pptx; None for PDFs, which are sent as documents.
    text: str | None
    #: Raw bytes for PDFs only.
    pdf_bytes: bytes | None = None

    @property
    def name(self) -> str:
        return self.path.name

    @property
    def is_pdf(self) -> bool:
        return self.pdf_bytes is not None


def ingest_sources(store: RunStore, settings: Settings) -> list[Source]:
    """Return sources that are new or changed since the last successful extraction.

    The index is only updated by `mark_processed`, called after reconcile succeeds — so a
    crash mid-extraction leaves the file pending rather than silently skipped next time.
    """
    # Pull the intake folder out of durable storage first. Without this the Cloud Run
    # container scans the empty `sources/` its own Dockerfile created and finds nothing, so
    # pushing a document to the repo was silently ignored — see materialize_sources().
    store.materialize_sources()

    sources_dir = settings.sources_dir
    if not sources_dir.is_dir():
        log.info("no sources/ directory")
        return []

    index = store.load_sources_index()
    pending: list[Source] = []

    for path in sorted(sources_dir.rglob("*")):
        if not path.is_file() or path.name.startswith("."):
            continue
        if path.name == "README.md":
            continue  # the folder's own instructions, not source material
        if path.suffix.lower() not in SUPPORTED:
            log.debug("skipping unsupported file %s", path.name)
            continue

        relative = path.relative_to(settings.repo_root).as_posix()
        digest = _sha256(path)
        if index.get(relative) == digest:
            continue

        try:
            source = _load(path, digest)
        except Exception as exc:  # noqa: BLE001 - one bad file must not fail the run
            log.warning("could not read %s: %s", path.name, exc)
            continue

        if source is not None:
            pending.append(source)
            log.info("%s %s", "changed:" if relative in index else "new:", relative)

    return pending


def mark_processed(sources: list[Source], store: RunStore, settings: Settings) -> None:
    """Record these sources as extracted. Call only after reconcile has succeeded."""
    index = store.load_sources_index()
    for source in sources:
        relative = source.path.relative_to(settings.repo_root).as_posix()
        index[relative] = source.sha256
    store.save_sources_index(index)
    log.info("recorded %d source(s) as processed", len(sources))


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        # Chunked so a large PDF doesn't sit in memory twice.
        while chunk := handle.read(1 << 20):
            digest.update(chunk)
    return digest.hexdigest()


def _load(path: Path, digest: str) -> Source | None:
    suffix = path.suffix.lower()

    if suffix in PDF_SUFFIXES:
        data = path.read_bytes()
        if len(data) > MAX_PDF_BYTES:
            log.warning(
                "%s is %.1f MB, over the %d MB limit — skipping",
                path.name,
                len(data) / 1024 / 1024,
                MAX_PDF_BYTES // 1024 // 1024,
            )
            return None
        return Source(path=path, sha256=digest, text=None, pdf_bytes=data)

    if suffix in TEXT_SUFFIXES:
        text = path.read_text(encoding="utf-8", errors="replace")
    elif suffix in DOCX_SUFFIXES:
        text = _docx_text(path)
    elif suffix in PPTX_SUFFIXES:
        text = _pptx_text(path)
    elif suffix in HTML_SUFFIXES:
        text = _html_text(path)
    else:  # pragma: no cover - guarded by the caller
        return None

    text = text.strip()
    if not text:
        log.warning("%s has no extractable text — skipping", path.name)
        return None
    if len(text) > MAX_CHARS:
        log.warning("%s is %d chars; truncating to %d", path.name, len(text), MAX_CHARS)
        text = text[:MAX_CHARS]

    return Source(path=path, sha256=digest, text=text)


def _docx_text(path: Path) -> str:
    """Paragraphs and table cells from a Word document.

    Table cells matter: a project report's deliverables are as likely to be in a table as
    in prose, and skipping them loses exactly the concrete detail a resume bullet needs.
    """
    from docx import Document

    document = Document(str(path))
    parts = [para.text for para in document.paragraphs if para.text.strip()]
    for table in document.tables:
        for row in table.rows:
            if cells := [cell.text.strip() for cell in row.cells if cell.text.strip()]:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


#: Tags whose *content* is code or styling, never prose. Their text must never reach the
#: model: a saved web page carries tens of kilobytes of minified JavaScript, and an LLM
#: handed that will happily invent a "project" out of variable names.
_HTML_SKIP_CONTENT = frozenset({"script", "style", "noscript", "template", "svg"})

#: Tags that end a line. Without these, `<p>A</p><p>B</p>` extracts as "AB" — two unrelated
#: sentences fused into one nonsense phrase, which is worse than losing either.
_HTML_BLOCK_TAGS = frozenset(
    {
        "address", "article", "aside", "blockquote", "br", "div", "dd", "dl", "dt",
        "fieldset", "figcaption", "figure", "footer", "form", "h1", "h2", "h3", "h4",
        "h5", "h6", "header", "hr", "li", "main", "nav", "ol", "p", "pre", "section",
        "table", "tbody", "thead", "title", "tr", "ul",
    }
)  # fmt: skip

#: Table cells are separated rather than newline-terminated, matching `_docx_text`: a row
#: reads as one record, which is how deliverables tables are usually written.
_HTML_CELL_TAGS = frozenset({"td", "th"})


class _TextExtractor(HTMLParser):
    """Collect readable text from HTML, dropping code and styling.

    Deliberately lenient: real files here are "Save as HTML" output from Word, Notion, or a
    browser, which is full of unbalanced tags, conditional comments and vendor namespaces
    (`<o:p>`). An unknown tag is simply not a block tag, so it is ignored rather than fatal.
    """

    def __init__(self) -> None:
        # convert_charrefs (the default) turns &amp; into & for us.
        super().__init__(convert_charrefs=True)
        self.parts: list[str] = []
        # A depth counter, not a boolean: nested <script> inside <script> is legal enough in
        # the wild, and a boolean would flip back on the first close and leak the remainder.
        self._skip_depth = 0

    def handle_starttag(self, tag: str, attrs: object) -> None:  # noqa: ARG002
        if tag in _HTML_SKIP_CONTENT:
            self._skip_depth += 1
        elif tag in _HTML_CELL_TAGS:
            self.parts.append(" | ")
        elif tag in _HTML_BLOCK_TAGS:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in _HTML_SKIP_CONTENT:
            self._skip_depth = max(0, self._skip_depth - 1)
        elif tag in _HTML_BLOCK_TAGS:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if self._skip_depth == 0:
            self.parts.append(data)


def _html_text(path: Path) -> str:
    """Readable text from an HTML file.

    Uses the standard library rather than an HTML library: the job is to strip markup, and
    `HTMLParser` does that without adding a dependency to a service whose image is already
    3.8GB. `<title>` is kept because it usually names the document.
    """
    parser = _TextExtractor()
    # errors="replace": a page saved as Big5 or with a mislabelled charset must not abort a
    # run. A few replacement characters cost one bullet; an exception costs the whole file.
    parser.feed(path.read_text(encoding="utf-8", errors="replace"))
    parser.close()

    # Collapse the whitespace the markup left behind: strip each line, drop empties, and
    # keep at most one blank line so paragraph structure survives for the model.
    lines = [line.strip(" \t|") for line in "".join(parser.parts).splitlines()]
    out: list[str] = []
    for line in lines:
        if line:
            out.append(" ".join(line.split()))
        elif out and out[-1]:
            out.append("")
    return "\n".join(out).strip()


def _pptx_text(path: Path) -> str:
    """Slide text, with each slide labelled.

    Slide numbers give the model structure to reason about — a closing deck's "results"
    slide reads differently from its agenda.
    """
    from pptx import Presentation

    presentation = Presentation(str(path))
    parts = []
    for index, slide in enumerate(presentation.slides, 1):
        lines = [
            shape.text.strip()
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text.strip()
        ]
        if lines:
            parts.append(f"--- Slide {index} ---\n" + "\n".join(lines))
    return "\n\n".join(parts)
